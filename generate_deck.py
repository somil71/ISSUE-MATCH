from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color constants
C_SURFACE0 = RGBColor(0x0a, 0x0c, 0x11)
C_SURFACE1 = RGBColor(0x12, 0x15, 0x1d)
C_SURFACE2 = RGBColor(0x1a, 0x1e, 0x2a)
C_SURFACE3 = RGBColor(0x23, 0x28, 0x38)
C_TEXT = RGBColor(0xb8, 0xbe, 0xc9)
C_TEXT_BRIGHT = RGBColor(0xf3, 0xf5, 0xf8)
C_ACCENT = RGBColor(0x63, 0x66, 0xf1)
C_SAFE = RGBColor(0x10, 0xb9, 0x81)
C_DANGER = RGBColor(0xf4, 0x3f, 0x5e)

# Layout
blank_slide_layout = prs.slide_layouts[6]

def set_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_SURFACE0

def add_title(slide, text):
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(12.333)
    height = Inches(1.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_BRIGHT

def add_text(slide, text, left, top, width, height, font_size=14, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return txBox

def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

# Slide 1: Title
s1 = prs.slides.add_slide(blank_slide_layout)
set_background(s1)
add_title(s1, "IssueMatch AI")
add_text(s1, "Find the right open-source issue. Know exactly why it's safe to touch, who to ask, and what to say.",
         Inches(1.5), Inches(1.8), Inches(10.333), Inches(0.8), font_size=20, align=PP_ALIGN.CENTER, color=C_TEXT_BRIGHT)

add_shape(s1, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.5), Inches(3.1), Inches(2), C_SURFACE1)
add_text(s1, "4", Inches(1.5), Inches(3.8), Inches(3.1), Inches(0.8), font_size=32, bold=True, align=PP_ALIGN.CENTER, font_name='Courier New', color=C_ACCENT)
add_text(s1, "AST Languages Parsed\n(Python, JS, TS, TSX)", Inches(1.6), Inches(4.5), Inches(2.9), Inches(0.8), font_size=12, align=PP_ALIGN.CENTER)

add_shape(s1, MSO_SHAPE.RECTANGLE, Inches(5.1), Inches(3.5), Inches(3.1), Inches(2), C_SURFACE1)
add_text(s1, "100%", Inches(5.1), Inches(3.8), Inches(3.1), Inches(0.8), font_size=32, bold=True, align=PP_ALIGN.CENTER, font_name='Courier New', color=C_ACCENT)
add_text(s1, "Deterministic Metrics\n(0 External LLM Calls)", Inches(5.2), Inches(4.5), Inches(2.9), Inches(0.8), font_size=12, align=PP_ALIGN.CENTER)

add_shape(s1, MSO_SHAPE.RECTANGLE, Inches(8.7), Inches(3.5), Inches(3.1), Inches(2), C_SURFACE1)
add_text(s1, "6", Inches(8.7), Inches(3.8), Inches(3.1), Inches(0.8), font_size=32, bold=True, align=PP_ALIGN.CENTER, font_name='Courier New', color=C_ACCENT)
add_text(s1, "Hops Max on Blast Map\n(Transitive Fan-in BFS)", Inches(8.8), Inches(4.5), Inches(2.9), Inches(0.8), font_size=12, align=PP_ALIGN.CENTER)

# Slide 2: Problem
s2 = prs.slides.add_slide(blank_slide_layout)
set_background(s2)
add_title(s2, "The Problem")

add_text(s2, "1. Misleading Tags", Inches(1), Inches(2), Inches(5), Inches(0.5), font_size=20, bold=True, color=C_TEXT_BRIGHT)
add_text(s2, "GitHub 'good first issue' labels routinely hide 'Here Be Dragons' codebase complexity.", Inches(1), Inches(2.6), Inches(5), Inches(1))

add_text(s2, "2. Unseen Blast Radius", Inches(1), Inches(3.8), Inches(5), Inches(0.5), font_size=20, bold=True, color=C_TEXT_BRIGHT)
add_text(s2, "Contributors start working blindly, unaware of deep transitive dependencies and high fan-in.", Inches(1), Inches(4.4), Inches(5), Inches(1))

add_text(s2, "3. Manual Context Transfer", Inches(1), Inches(5.6), Inches(5), Inches(0.5), font_size=20, bold=True, color=C_TEXT_BRIGHT)
add_text(s2, "Maintainers waste hours answering 'where do I add the tests?' and 'who wrote this?'", Inches(1), Inches(6.2), Inches(5), Inches(1))

add_shape(s2, MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2), Inches(6), Inches(4.8), C_SURFACE1)
add_text(s2, "EXPECTED", Inches(7), Inches(2.5), Inches(2.5), Inches(0.5), font_size=14, bold=True, color=C_SAFE, font_name='Courier New')
add_shape(s2, MSO_SHAPE.OVAL, Inches(7.5), Inches(3), Inches(0.8), Inches(0.8), C_SURFACE2)
add_shape(s2, MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(3.3), Inches(0.6), Inches(0.2), C_SURFACE3)
add_shape(s2, MSO_SHAPE.OVAL, Inches(9.3), Inches(3), Inches(0.8), Inches(0.8), C_SURFACE2)

add_text(s2, "REALITY", Inches(7), Inches(4.3), Inches(2.5), Inches(0.5), font_size=14, bold=True, color=C_DANGER, font_name='Courier New')
add_shape(s2, MSO_SHAPE.OVAL, Inches(7.5), Inches(4.8), Inches(0.8), Inches(0.8), C_SURFACE2)
for i in range(3):
    add_shape(s2, MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(4.2 + i*0.45), Inches(0.6), Inches(0.2), C_SURFACE3)
    add_shape(s2, MSO_SHAPE.OVAL, Inches(9.3), Inches(4.0 + i*0.45), Inches(0.6), Inches(0.6), C_SURFACE2)
add_shape(s2, MSO_SHAPE.RIGHT_ARROW, Inches(10.1), Inches(4.2), Inches(0.6), Inches(0.2), C_SURFACE3)
add_shape(s2, MSO_SHAPE.OVAL, Inches(10.9), Inches(4.0), Inches(0.6), Inches(0.6), C_SURFACE2)

# Slide 3: Market & Timing
s3 = prs.slides.add_slide(blank_slide_layout)
set_background(s3)
add_title(s3, "Why This, Why Now")

add_text(s3, "The Gap: AI Guesses vs. Static Proof", Inches(1), Inches(2), Inches(11), Inches(0.5), font_size=22, bold=True, color=C_TEXT_BRIGHT)
add_text(s3, "Current AI coding tools rely on LLM hallucinations for risk assessment. We use pure, deterministic tree-sitter AST parsing to mathematically prove the risk before a single line is written.", Inches(1), Inches(2.6), Inches(11), Inches(0.8), font_size=15)

add_shape(s3, MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.8), Inches(5.3), Inches(2.8), C_SURFACE1)
add_text(s3, "LLM Approach", Inches(1.3), Inches(4), Inches(4.7), Inches(0.5), font_size=18, bold=True, color=C_DANGER)
add_text(s3, "1. Opaque reasoning\n2. Prone to hallucinations\n3. Unbounded API latency\n4. Expensive token costs", Inches(1.3), Inches(4.5), Inches(4.7), Inches(1.8), font_size=14)

add_shape(s3, MSO_SHAPE.RECTANGLE, Inches(7), Inches(3.8), Inches(5.3), Inches(2.8), C_SURFACE1)
add_text(s3, "IssueMatch Approach", Inches(7.3), Inches(4), Inches(4.7), Inches(0.5), font_size=18, bold=True, color=C_SAFE)
add_text(s3, "1. 100% Deterministic AST\n2. 0 External API Calls\n3. Fast Local Caching (Redis)\n4. Transparent Formula Weights", Inches(7.3), Inches(4.5), Inches(4.7), Inches(1.8), font_size=14)

# Slide 4: The Idea
s4 = prs.slides.add_slide(blank_slide_layout)
set_background(s4)
add_title(s4, "The Three Pillars")

for i, title, desc in [
    (0, "Static Risk Proof", "Deterministic blast radius from Fan-In (0.35), Complexity (0.30), Git Churn (0.15), and Test Proximity (0.20)."),
    (1, "Local Skill Match", "Cosine similarity ranking of issue text vs user skills using local sentence-transformers and TF-IDF terms."),
    (2, "Playbook Gen", "Auto-drafts PR descriptions and branch names using real Git facts, bot-filtered commit authors, and test locs.")
]:
    left = Inches(1 + i*3.9)
    top = Inches(2.5)
    add_shape(s4, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.6), Inches(3.5), C_SURFACE1)
    add_text(s4, f"Pillar {i+1}", left + Inches(0.3), top + Inches(0.4), Inches(3), Inches(0.4), font_size=14, bold=True, color=C_ACCENT, font_name='Courier New')
    add_text(s4, title, left + Inches(0.3), top + Inches(0.9), Inches(3), Inches(0.4), font_size=20, bold=True, color=C_TEXT_BRIGHT)
    add_text(s4, desc, left + Inches(0.3), top + Inches(1.5), Inches(3), Inches(1.8), font_size=14)

# Slide 5: How It Works
s5 = prs.slides.add_slide(blank_slide_layout)
set_background(s5)
add_title(s5, "System Architecture: Pure Determinism")

add_shape(s5, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(12.333), Inches(1.5), C_SURFACE1)
add_text(s5, "INPUT", Inches(0.6), Inches(2.1), Inches(1.5), Inches(0.3), font_size=12, color=C_TEXT_BRIGHT, font_name='Courier New')

add_shape(s5, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.7), Inches(12.333), Inches(1.5), C_SURFACE2)
add_text(s5, "ANALYSIS", Inches(0.6), Inches(3.8), Inches(1.5), Inches(0.3), font_size=12, color=C_TEXT_BRIGHT, font_name='Courier New')

add_shape(s5, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.4), Inches(12.333), Inches(1.5), C_SURFACE3)
add_text(s5, "OUTPUT", Inches(0.6), Inches(5.5), Inches(1.5), Inches(0.3), font_size=12, color=C_TEXT_BRIGHT, font_name='Courier New')

def node(slide, x, y, w, h, text, font_size=10):
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h), C_SURFACE0)
    add_text(slide, text, Inches(x+0.1), Inches(y+0.1), Inches(w-0.2), Inches(h-0.2), font_size=font_size, align=PP_ALIGN.CENTER, font_name='Courier New')

node(s5, 1.5, 2.35, 2.5, 0.8, "Git Clone\n(--filter=blob:none)")
node(s5, 5, 2.35, 3, 0.8, "Repo Metadata\n(GitHub API)")
node(s5, 9, 2.35, 3, 0.8, "Contributor Skills\n(Auth/Profile)")

node(s5, 1, 4.05, 2.5, 0.8, "Tree-Sitter AST\nFunction + Call Edges")
node(s5, 4, 4.05, 2.7, 0.8, "Static Fact Extraction\n(Fan-in, Complexity, Git)")
node(s5, 7.2, 4.05, 2.5, 0.8, "Formula Weights\n(0.35/0.30/0.20/0.15)")
node(s5, 10.2, 4.05, 2.5, 0.8, "Local Embeddings\nTF-IDF & sentence-trs")

node(s5, 1.5, 5.75, 3, 0.8, "Blast Radius Score\nStart Here vs Dragons")
node(s5, 5.2, 5.75, 3, 0.8, "First Merge Path\n(3 lowest risk functions)")
node(s5, 8.9, 5.75, 3, 0.8, "Actionable Playbook\nPR Drafts & X-Refs")

add_shape(s5, MSO_SHAPE.RIGHT_ARROW, Inches(2.75), Inches(3.15), Inches(0.2), Inches(0.9), C_SURFACE1)
add_shape(s5, MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(3.15), Inches(0.2), Inches(0.9), C_SURFACE1)
add_shape(s5, MSO_SHAPE.RIGHT_ARROW, Inches(10.5), Inches(3.15), Inches(0.2), Inches(0.9), C_SURFACE1)
add_shape(s5, MSO_SHAPE.RIGHT_ARROW, Inches(3.5), Inches(4.35), Inches(0.5), Inches(0.2), C_SURFACE1)
add_shape(s5, MSO_SHAPE.RIGHT_ARROW, Inches(6.7), Inches(4.35), Inches(0.5), Inches(0.2), C_SURFACE1)
add_shape(s5, MSO_SHAPE.RIGHT_ARROW, Inches(9.7), Inches(4.35), Inches(0.5), Inches(0.2), C_SURFACE1)

# Slide 6: Tech Stack
s6 = prs.slides.add_slide(blank_slide_layout)
set_background(s6)
add_title(s6, "The Stack")

add_shape(s6, MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(11.333), Inches(1.3), C_SURFACE1)
add_text(s6, "Frontend", Inches(1.3), Inches(2.4), Inches(2), Inches(0.5), font_size=18, bold=True, color=C_TEXT_BRIGHT)
add_text(s6, "React 19 | Vite | TypeScript\nTailwind CSS v4 | React Query", Inches(3.5), Inches(2.2), Inches(8), Inches(0.9), font_size=16, font_name='Courier New', color=C_ACCENT)

add_shape(s6, MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.8), Inches(11.333), Inches(1.3), C_SURFACE2)
add_text(s6, "Backend", Inches(1.3), Inches(4.2), Inches(2), Inches(0.5), font_size=18, bold=True, color=C_TEXT_BRIGHT)
add_text(s6, "FastAPI | tree-sitter | sentence-transformers\nscikit-learn | spaCy", Inches(3.5), Inches(4.0), Inches(8.5), Inches(0.9), font_size=15, font_name='Courier New', color=C_ACCENT)

add_shape(s6, MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.6), Inches(11.333), Inches(1.3), C_SURFACE3)
add_text(s6, "Data & Ops", Inches(1.3), Inches(6.0), Inches(2), Inches(0.5), font_size=18, bold=True, color=C_TEXT_BRIGHT)
add_text(s6, "PostgreSQL | Redis (1-hour TTL cache)\nDocker Compose", Inches(3.5), Inches(5.8), Inches(8), Inches(0.9), font_size=16, font_name='Courier New', color=C_ACCENT)

# Slide 7: Product Walkthrough
s7 = prs.slides.add_slide(blank_slide_layout)
set_background(s7)
add_title(s7, "Inside the Engine")

add_shape(s7, MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(11.333), Inches(5.2), C_SURFACE1)
add_shape(s7, MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(11.333), Inches(0.6), C_SURFACE2)
add_text(s7, "api.github.com/repos/owner/name", Inches(1.2), Inches(1.95), Inches(5), Inches(0.3), font_size=12, font_name='Courier New', color=C_TEXT)

add_shape(s7, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.8), Inches(4.5), Inches(1.2), C_SURFACE2)
add_text(s7, "auth_middleware()", Inches(1.7), Inches(3.0), Inches(4), Inches(0.3), font_size=14, bold=True, color=C_TEXT_BRIGHT, font_name='Courier New')
add_text(s7, "Blast Radius: START HERE", Inches(1.7), Inches(3.4), Inches(4), Inches(0.3), font_size=12, color=C_SAFE, bold=True)
add_text(s7, "2 callers, no branches, has a nearby test.", Inches(1.7), Inches(3.7), Inches(4), Inches(0.3), font_size=10)

add_shape(s7, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.3), Inches(4.5), Inches(1.2), C_SURFACE2)
add_text(s7, "execute_query_unsafe()", Inches(1.7), Inches(4.5), Inches(4), Inches(0.3), font_size=14, bold=True, color=C_TEXT_BRIGHT, font_name='Courier New')
add_text(s7, "Blast Radius: HERE BE DRAGONS", Inches(1.7), Inches(4.9), Inches(4), Inches(0.3), font_size=12, color=C_DANGER, bold=True)
add_text(s7, "45 callers, 12 branches, no tests nearby.", Inches(1.7), Inches(5.2), Inches(4), Inches(0.3), font_size=10)

add_shape(s7, MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2.8), Inches(5.3), Inches(2.7), C_SURFACE0)
add_text(s7, "Trust Panel Network Log", Inches(6.7), Inches(3), Inches(4.8), Inches(0.3), font_size=12, bold=True, color=C_TEXT_BRIGHT)
add_text(s7, "Host: api.github.com", Inches(6.7), Inches(3.4), Inches(4.8), Inches(0.3), font_size=12, font_name='Courier New', color=C_SAFE)
add_text(s7, "Status: VERIFIED (0 EXTERNAL)", Inches(6.7), Inches(3.7), Inches(4.8), Inches(0.3), font_size=12, font_name='Courier New', color=C_SAFE)
add_text(s7, "-----------------------------", Inches(6.7), Inches(4.0), Inches(4.8), Inches(0.3), font_size=12, font_name='Courier New')
add_text(s7, "Weights Array:", Inches(6.7), Inches(4.3), Inches(4.8), Inches(0.3), font_size=12, font_name='Courier New')
add_text(s7, "  fan_in: 0.35\n  complexity: 0.30\n  test_prox: 0.20\n  git_churn: 0.15", Inches(6.7), Inches(4.6), Inches(4.8), Inches(1), font_size=12, font_name='Courier New')

# Slide 8: Why it Wins + Close
s8 = prs.slides.add_slide(blank_slide_layout)
set_background(s8)
add_title(s8, "Why We Win")

add_shape(s8, MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.2), Inches(5.2), Inches(1.5), C_SURFACE1)
add_text(s8, "1. Structural Privacy", Inches(1.3), Inches(2.4), Inches(4.6), Inches(0.4), font_size=16, bold=True, color=C_TEXT_BRIGHT)
add_text(s8, "Only ONE logged outbound HTTP route (GitHub). Falsifiable by codebase inspection.", Inches(1.3), Inches(2.9), Inches(4.6), Inches(0.7), font_size=13)

add_shape(s8, MSO_SHAPE.RECTANGLE, Inches(7.133), Inches(2.2), Inches(5.2), Inches(1.5), C_SURFACE1)
add_text(s8, "2. Real Code Context", Inches(7.433), Inches(2.4), Inches(4.6), Inches(0.4), font_size=16, bold=True, color=C_TEXT_BRIGHT)
add_text(s8, "Bot commits are stripped. Real human code owners are extracted and tagged.", Inches(7.433), Inches(2.9), Inches(4.6), Inches(0.7), font_size=13)

add_shape(s8, MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.0), Inches(5.2), Inches(1.5), C_SURFACE1)
add_text(s8, "3. True Determinism", Inches(1.3), Inches(4.2), Inches(4.6), Inches(0.4), font_size=16, bold=True, color=C_TEXT_BRIGHT)
add_text(s8, "Not a single LLM call is made. Everything is calculated from static AST and Git history.", Inches(1.3), Inches(4.7), Inches(4.6), Inches(0.7), font_size=13)

add_shape(s8, MSO_SHAPE.RECTANGLE, Inches(7.133), Inches(4.0), Inches(5.2), Inches(1.5), C_SURFACE1)
add_text(s8, "4. Fast Caching", Inches(7.433), Inches(4.2), Inches(4.6), Inches(0.4), font_size=16, bold=True, color=C_TEXT_BRIGHT)
add_text(s8, "Redis caches graph analyses (1-hr TTL). Subsequent queries are lightning fast.", Inches(7.433), Inches(4.7), Inches(4.6), Inches(0.7), font_size=13)

add_text(s8, "Open Source needs proof, not promises.", Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.8), font_size=32, bold=True, align=PP_ALIGN.CENTER, color=C_ACCENT)

prs.save("pitch_deck.pptx")
